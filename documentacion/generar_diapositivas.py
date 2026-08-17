import subprocess
import sys
import os

# Asegurar que la librería python-pptx esté instalada
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Instalando 'python-pptx' necesario para generar la presentación...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def crear_presentacion():
    prs = Presentation()
    
    # Configurar dimensiones a Widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Paleta de Colores
    TEAL = RGBColor(16, 152, 173)       # Color Primario
    DARK = RGBColor(33, 37, 41)         # Texto oscuro
    MUTED = RGBColor(134, 142, 150)     # Gris
    ORANGE = RGBColor(255, 146, 43)     # Acento
    WHITE = RGBColor(255, 255, 255)
    
    # Usar diseño en blanco (Slide layout index 6 es blanco)
    blank_layout = prs.slide_layouts[6]
    
    # Helper para agregar un fondo de color sólido (opcional)
    def pintar_fondo(slide, rgb_color):
      background = slide.background
      fill = background.fill
      fill.solid()
      fill.fore_color.rgb = rgb_color

    # Helper para crear títulos comunes
    def agregar_titulo(slide, texto, color=TEAL):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = texto
        p.font.name = 'Arial'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        return title_box

    # ==========================================
    # DIAPOSITIVA 1: PORTADA (Fondo Oscuro)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    pintar_fondo(slide1, DARK)
    
    # Cuadro de texto para portada
    portada_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
    tf1 = portada_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "PROYECTO DE INNOVACIÓN TECNOLÓGICA"
    p1.font.name = 'Arial'
    p1.font.size = Pt(20)
    p1.font.color.rgb = ORANGE
    p1.font.bold = True
    
    p2 = tf1.add_paragraph()
    p2.text = "EncuentraMiMascota"
    p2.font.name = 'Arial'
    p2.font.size = Pt(54)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(20)
    
    p3 = tf1.add_paragraph()
    p3.text = "Plataforma Web Responsive para el Reporte y Búsqueda de Mascotas Perdidas"
    p3.font.name = 'Arial'
    p3.font.size = Pt(22)
    p3.font.color.rgb = TEAL
    p3.space_after = Pt(40)
    
    p4 = tf1.add_paragraph()
    p4.text = "Estudiante: [Tu Nombre Completo]\nCurso: 6to de Secundaria BTH - Sistemas Informáticos\nUnidad Educativa: [Nombre de tu Colegio]"
    p4.font.name = 'Arial'
    p4.font.size = Pt(16)
    p4.font.color.rgb = WHITE

    # ==========================================
    # DIAPOSITIVA 2: INTRODUCCIÓN Y DIAGNÓSTICO
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    agregar_titulo(slide2, "Introducción y Diagnóstico")
    
    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "La Realidad en Nuestra Comunidad:"
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.space_after = Pt(15)
    
    bullets = [
        "El extravío de mascotas genera alta angustia y dispersión de la información.",
        "Los métodos tradicionales (afiches, publicaciones sueltas en redes sociales) tienen bajo alcance y expiran rápidamente sin control.",
        "Según nuestro sondeo local, el 80% considera ineficiente el uso de grupos genéricos en Facebook o WhatsApp para búsquedas rápidas.",
        "Solución Propuesta: Centralizar los reportes en un sitio web optimizado, interactivo y con filtros de búsqueda instantáneos."
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(20)
        p_b.font.color.rgb = MUTED
        p_b.space_after = Pt(10)

    # ==========================================
    # DIAPOSITIVA 3: FORMULACIÓN Y OBJETIVOS
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    agregar_titulo(slide3, "Formulación y Objetivos")
    
    content_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf3 = content_box3.text_frame
    tf3.word_wrap = True
    
    p = tf3.paragraphs[0]
    p.text = "¿Pregunta Clave del Problema?:"
    p.font.name = 'Arial'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(5)
    
    p_q = tf3.add_paragraph()
    p_q.text = '¿De qué manera el desarrollo de una aplicación web responsive basada en una arquitectura cliente-servidor y base de datos relacional puede optimizar el tiempo de búsqueda y reporte de mascotas perdidas en nuestra comunidad?'
    p_q.font.name = 'Arial'
    p_q.font.size = Pt(18)
    p_q.font.italic = True
    p_q.font.color.rgb = DARK
    p_q.space_after = Pt(25)
    
    p_obj = tf3.add_paragraph()
    p_obj.text = "Objetivo General:"
    p_obj.font.name = 'Arial'
    p_obj.font.size = Pt(22)
    p_obj.font.bold = True
    p_obj.font.color.rgb = TEAL
    p_obj.space_after = Pt(5)
    
    p_obj_val = tf3.add_paragraph()
    p_obj_val.text = 'Desarrollar una aplicación web responsive utilizando HTML5, CSS3, JavaScript Vanilla en el frontend y Node.js con SQLite en el backend para reportar, buscar y comentar mascotas perdidas.'
    p_obj_val.font.name = 'Arial'
    p_obj_val.font.size = Pt(18)
    p_obj_val.font.color.rgb = MUTED
    p_obj_val.space_after = Pt(15)

    # ==========================================
    # DIAPOSITIVA 4: MARCO REFERENCIAL (TECNOLOGÍAS)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    agregar_titulo(slide4, "Marco Referencial y Tecnológico")
    
    content_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf4 = content_box4.text_frame
    tf4.word_wrap = True
    
    techs = [
        ("Frontend: HTML5, CSS3 y JS Vanilla", "Código nativo, sin frameworks complejos, optimizado para celulares (Responsive Design) usando CSS Grid y Flexbox."),
        ("Backend: Node.js con Express", "Servidor ligero encargado de exponer una API REST segura para comunicarse con el cliente."),
        ("Base de Datos: SQLite", "Motor de base de datos relacional liviano autocontenido en un solo archivo local ('mascotas.db')."),
        ("Seguridad: Encriptación con Bcrypt y JWT", "Contraseñas encriptadas de forma segura y control de sesión del usuario mediante Tokens web.")
    ]
    
    for titulo, desc in techs:
        p_t = tf4.add_paragraph()
        p_t.text = titulo
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = DARK
        
        p_d = tf4.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(16)
        p_d.font.color.rgb = MUTED
        p_d.space_after = Pt(15)

    # ==========================================
    # DIAPOSITIVA 5: CARACTERÍSTICAS DEL PRODUCTO
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    agregar_titulo(slide5, "Características de EncuentraMiMascota")
    
    content_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf5 = content_box5.text_frame
    tf5.word_wrap = True
    
    bullets5 = [
        "Registro y Edición de Anuncios: Nombre, Ubicación de extravío, Teléfono, Recompensa (en Bs.), Descripción detallada y carga de Foto real.",
        "Actualización Dinámica de Fotos: Los dueños y administradores pueden modificar anuncios y cambiar fotografías con previsualización instantánea.",
        "Gestión Limpia de Servidor: Al reemplazar fotos o borrar anuncios, el backend limpia automáticamente los archivos obsoletos del disco.",
        "Cambio de Estados: Los reportes cuentan con los estados 'Perdido' (color rojo) y 'Encontrado' (color verde) para un mejor control visual.",
        "Sección de Comentarios: Los usuarios con sesión iniciada pueden aportar pistas y debatir en los anuncios para cooperar.",
        "Buscador Inteligente: Entrada de búsqueda interactiva que filtra por nombre, zona o descripción en tiempo real."
    ]
    for b in bullets5:
        p_b = tf5.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(18)
        p_b.font.color.rgb = MUTED
        p_b.space_after = Pt(15)

    # ==========================================
    # DIAPOSITIVA 6: ARQUITECTURA DE ROLES
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    agregar_titulo(slide6, "Arquitectura y Control de Roles")
    
    content_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf6 = content_box6.text_frame
    tf6.word_wrap = True
    
    roles = [
        ("Invitado (Sin inicio de sesión):", "Puede entrar al portal, ver todos los anuncios de mascotas perdidas, utilizar el buscador y leer comentarios. Rol meramente observador."),
        ("Usuario Registrado (Sesión activa):", "Puede crear publicaciones de mascotas perdidas, marcarlas como encontradas o eliminarlas (solo sus propias publicaciones). Puede comentar en cualquier publicación."),
        ("Administrador (Control Total):", "Acceso al panel de administración protegido. Puede crear/editar/eliminar cualquier usuario registrado y tiene el poder de modificar o eliminar anuncios inapropiados de cualquier usuario.")
    ]
    
    for titulo, desc in roles:
        p_t = tf6.add_paragraph()
        p_t.text = titulo
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = TEAL
        
        p_d = tf6.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(16)
        p_d.font.color.rgb = MUTED
        p_d.space_after = Pt(15)

    # ==========================================
    # DIAPOSITIVA 7: METODOLOGÍA Y RESULTADOS
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    agregar_titulo(slide7, "Metodología y Resultados")
    
    content_box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf7 = content_box7.text_frame
    tf7.word_wrap = True
    
    p = tf7.paragraphs[0]
    p.text = "Metodología Aplicada:"
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.space_after = Pt(5)
    
    p_met = tf7.add_paragraph()
    p_met.text = "• Investigación-Acción Participativa (IAP): Identificación directa de la necesidad, desarrollo de prototipo, y realización de encuestas y pruebas de usabilidad con vecinos locales."
    p_met.font.name = 'Arial'
    p_met.font.size = Pt(16)
    p_met.font.color.rgb = MUTED
    p_met.space_after = Pt(15)

    p_res = tf7.add_paragraph()
    p_res.text = "Resultados Clave Obtenidos:"
    p_res.font.name = 'Arial'
    p_res.font.size = Pt(20)
    p_res.font.bold = True
    p_res.font.color.rgb = DARK
    p_res.space_after = Pt(5)
    
    p_res_val = tf7.add_paragraph()
    p_res_val.text = "• Centralización efectiva: Los reportes están ordenados y no se pierden en hilos temporales.\n• Tiempos de carga mínimos: Gracias al uso de SQLite local y Node.js.\n• Interfaz responsiva probada con éxito en dispositivos móviles."
    p_res_val.font.name = 'Arial'
    p_res_val.font.size = Pt(16)
    p_res_val.font.color.rgb = MUTED
    p_res_val.space_after = Pt(15)

    # ==========================================
    # DIAPOSITIVA 8: RECURSOS Y COSTOS
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    agregar_titulo(slide8, "Recursos y Cálculo de Costos")
    
    content_box8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf8 = content_box8.text_frame
    tf8.word_wrap = True
    
    costos = [
        ("Recursos Materiales y Software:", "Computadora portátil, Visual Studio Code, Git, Node.js y SQLite. Inversión inicial en software de desarrollo: 0 Bs. (Herramientas libres)."),
        ("Costos de Operación Estimados:", "Energía eléctrica y conexión de banda ancha a internet: 170 Bs. mensuales. Plan gratuito de despliegue en la nube."),
        ("Costo Estimado de Desarrollo (Mano de obra):", "120 horas de programación valoradas en 20 Bs. / hora. Costo simulado del proyecto de desarrollo: 2,400 Bs.")
    ]
    
    for titulo, desc in costos:
        p_t = tf8.add_paragraph()
        p_t.text = titulo
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = ORANGE
        
        p_d = tf8.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(16)
        p_d.font.color.rgb = MUTED
        p_d.space_after = Pt(15)

    # ==========================================
    # DIAPOSITIVA 9: ESTRATEGIA DE MEJORA Y CONCLUSIONES
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    agregar_titulo(slide9, "Estrategia de Mejora y Conclusiones")
    
    content_box9 = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf9 = content_box9.text_frame
    tf9.word_wrap = True
    
    p = tf9.paragraphs[0]
    p.text = "Estrategias de Mejora Futura:"
    p.font.name = 'Arial'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK
    
    p_m = tf9.add_paragraph()
    p_m.text = "1. Integración de Mapas Interactivos (API Leaflet/Google Maps).\n2. Notificaciones push al móvil convirtiendo la web en una PWA."
    p_m.font.name = 'Arial'
    p_m.font.size = Pt(16)
    p_m.font.color.rgb = MUTED
    p_m.space_after = Pt(20)

    p_c = tf9.add_paragraph()
    p_c.text = "Conclusiones de Grado:"
    p_c.font.name = 'Arial'
    p_c.font.size = Pt(18)
    p_c.font.bold = True
    p_c.font.color.rgb = TEAL
    
    p_c_val = tf9.add_paragraph()
    p_c_val.text = "• Se completaron todos los requisitos funcionales del proyecto bajo un enfoque ordenado y limpio.\n• El Bachillerato Técnico (BTH) nos habilita para diseñar e implementar software que resuelva problemas reales de la comunidad."
    p_c_val.font.name = 'Arial'
    p_c_val.font.size = Pt(16)
    p_c_val.font.color.rgb = MUTED

    # ==========================================
    # DIAPOSITIVA 10: CIERRE (Fondo Oscuro)
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    pintar_fondo(slide10, DARK)
    
    cierre_box = slide10.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(3.0))
    tf10 = cierre_box.text_frame
    tf10.word_wrap = True
    
    p_thx = tf10.paragraphs[0]
    p_thx.text = "Muchas Gracias por su Atención"
    p_thx.alignment = PP_ALIGN.CENTER
    p_thx.font.name = 'Arial'
    p_thx.font.size = Pt(44)
    p_thx.font.bold = True
    p_thx.font.color.rgb = WHITE
    p_thx.space_after = Pt(20)
    
    p_ask = tf10.add_paragraph()
    p_ask.text = "¿Preguntas del Tribunal Evaluador?"
    p_ask.alignment = PP_ALIGN.CENTER
    p_ask.font.name = 'Arial'
    p_ask.font.size = Pt(28)
    p_ask.font.color.rgb = ORANGE
    p_ask.font.bold = True

    # Guardar presentación en disco
    dir_path = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(dir_path, "presentacion_defensa.pptx")
    prs.save(output_path)
    print(f"Presentación PPTX guardada exitosamente en:\n{output_path}")

if __name__ == "__main__":
    crear_presentacion()
